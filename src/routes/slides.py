"""
RepoLM — PPTX slide generation endpoint.
"""

import os
import re
import uuid

from fastapi import APIRouter, Request
from fastapi.responses import JSONResponse, FileResponse

from config import OUTPUT_DIR

router = APIRouter()


@router.post("/api/slides-pptx")
async def slides_pptx(request: Request):
    """Generate a .pptx file from slide markdown."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return JSONResponse({"error": "python-pptx not installed"}, 500)

    body = await request.json()
    markdown = body.get("markdown", "")
    repo_name = body.get("repo_name", "Repository")

    raw_slides = re.split(r'\n---+\n', markdown)
    slides_data = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        lines = raw.split('\n')
        title = ""
        bullets = []
        code_block = []
        in_code = False
        key_takeaway = ""
        for line in lines:
            if line.startswith('```'):
                in_code = not in_code
                continue
            if in_code:
                code_block.append(line)
                continue
            if line.startswith('# '):
                title = line[2:].strip()
            elif line.startswith('**Key Takeaway:**'):
                key_takeaway = line.replace('**Key Takeaway:**', '').replace('**', '').strip()
            elif line.startswith('- ') or line.startswith('* '):
                bullets.append(line[2:].strip())
            elif line.strip() and not title:
                title = line.strip().lstrip('#').strip()
        slides_data.append({"title": title, "bullets": bullets, "code": code_block, "takeaway": key_takeaway})

    if not slides_data:
        return JSONResponse({"error": "No slides found in markdown"}, 400)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(0x09, 0x09, 0x0b)
    title_color = RGBColor(0xa7, 0x8b, 0xfa)
    text_color = RGBColor(0xe5, 0xe7, 0xeb)
    accent_color = RGBColor(0x8b, 0x5c, 0xf6)

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        left, top = Inches(0.8), Inches(0.5)
        width, height = Inches(11.5), Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd["title"] or f"Slide {i+1}"
        p.font.size = Pt(36)
        p.font.color.rgb = title_color
        p.font.bold = True

        if sd["bullets"]:
            top_b = Inches(2.0)
            txBox2 = slide.shapes.add_textbox(Inches(1.0), top_b, Inches(11), Inches(3.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for j, bullet in enumerate(sd["bullets"]):
                p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                p2.text = "•  " + bullet
                p2.font.size = Pt(20)
                p2.font.color.rgb = text_color
                p2.space_after = Pt(12)

        if sd["code"]:
            code_top = Inches(5.5) if sd["bullets"] else Inches(2.0)
            code_text = "\n".join(sd["code"])
            txBox3 = slide.shapes.add_textbox(Inches(1.0), code_top, Inches(11), Inches(1.5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = code_text
            p3.font.size = Pt(14)
            p3.font.color.rgb = RGBColor(0xc0, 0xc0, 0xc0)
            p3.font.name = "Consolas"

        if sd["takeaway"]:
            txBox4 = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8))
            tf4 = txBox4.text_frame
            tf4.word_wrap = True
            p4 = tf4.paragraphs[0]
            p4.text = "⚡ " + sd["takeaway"]
            p4.font.size = Pt(16)
            p4.font.color.rgb = accent_color
            p4.font.italic = True

    pptx_path = os.path.join(OUTPUT_DIR, f"slides_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(pptx_path)
    return FileResponse(pptx_path,
                       media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                       filename=f"{repo_name}_slides.pptx")
